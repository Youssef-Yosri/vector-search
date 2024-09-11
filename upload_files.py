def upload():
    import os
    import glob
    import mimetypes
    import PyPDF2
    from docx import Document
    from pptx import Presentation
    from transformers import BertModel, BertTokenizer
    import torch
    import numpy as np
    import re
    from pymongo import MongoClient
    from datetime import datetime
    from pymongo import MongoClient
    import numpy as np
    import faiss
    
    folder_path = input("Enter folder path to upload files from (test with the Samples folder if you have no folder): ")
    files = glob.glob(os.path.join(folder_path, '*'))
    
    for file in files:
        print(file)

    def extract_text_from_pdf(file_path):
        text = ""
        with open(file_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            for page in range(len(reader.pages)):
                text += reader.pages[page].extract_text()
        return text

    def extract_text_from_docx(file_path):
        doc = Document(file_path)
        text = '\n'.join([para.text for para in doc.paragraphs])
        return text

    def extract_text_from_pptx(file_path):
        prs = Presentation(file_path)
        text = ''
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + '\n'
        return text

    def extract_text_from_file(file_path):
        mime_type, _ = mimetypes.guess_type(file_path)
        
        if mime_type == 'application/pdf':
            return extract_text_from_pdf(file_path)
        elif mime_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
            return extract_text_from_docx(file_path)
        elif mime_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
            return extract_text_from_pptx(file_path)
        else:
            raise ValueError("Unsupported file type")

    def split_text_into_chunks(text, max_length=512):
        paragraphs = text.split('\n\n')
        chunks = []
        current_chunk = ""
    
        for paragraph in paragraphs:
            if len(current_chunk) + len(paragraph) <= max_length:
                current_chunk += paragraph + " "
            else:
                chunks.append(current_chunk.strip())
                current_chunk = paragraph + " "
    
        if current_chunk:
            chunks.append(current_chunk.strip())

        chunks = [chunk for chunk in chunks if chunk.strip()]
    
        return chunks

    tokenizer = BertTokenizer.from_pretrained('bert-large-uncased')
    model = BertModel.from_pretrained('bert-large-uncased')
    
    def embed_text(text):
        inputs = tokenizer(text, return_tensors='pt', truncation=True, padding=True)
        with torch.no_grad():
            embedding = model(**inputs).last_hidden_state.mean(dim=1)
        return embedding
    
    number = 0
    
    for file_path in files:
        try:
            text_content = extract_text_from_file(file_path)
    
            email_pattern = r'[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}'
            emails = re.findall(email_pattern, text_content)
    
            if emails:
                email = emails[0]
            else:
                email = "example@email.com"
    
            chunks = split_text_into_chunks(text_content)
            
            chunk_embeddings = [embed_text(chunk) for chunk in chunks]
            
            def average_embeddings(embeddings):
                return torch.mean(torch.stack(embeddings), dim=0)
            
            document_embedding = average_embeddings(chunk_embeddings)

            embedding_list = document_embedding.tolist()

            client = MongoClient('mongodb://localhost:27017/')
            db = client['VectorDBPython']
            collection = db['CVs']

            document = {
                "metadata": {
                    "file_path": file_path,
                    "email": email,
                    "created_at": datetime.now().isoformat(),
                    "tags": ["tag1", "tag2"],
                },
                "embedding": embedding_list,
                "text_excerpt": text_content,
            }

            result = collection.insert_one(document)
            print(f"Inserted document ID: {result.inserted_id}")
            number = number + 1
    
        except ValueError as e:
            print(e)
    print(f"{number} documents added")

    connectionString = input("Enter your connection string (the default local host connection string of mongodb is mongodb://localhost:27017/): ")
    dbName = input("Enter your database name: ")
    collectionName = input("Enter your collection name: ")
    
    client = MongoClient(connectionString)
    db = client[dbName]
    collection = db[collectionName]
    
    vectors = []
    text = []
    for doc in collection.find():
        vectors.append(doc['embedding'])
        text.append(doc['text_excerpt']) 
    
    vectors = np.array(vectors, dtype='float32')
    
    vectors = np.squeeze(vectors)
    print("Vectors shape:", vectors.shape)
    
    norms = np.linalg.norm(vectors, axis=1, keepdims=True)
    vectors /= norms
    
    print("Norms statistics - Min:", norms.min(), "Max:", norms.max(), "Mean:", norms.mean())
    
    norms = np.linalg.norm(vectors, axis=1, keepdims=True)
    vectors /= norms
    print("Recomputed norms - Min:", norms.min(), "Max:", norms.max(), "Mean:", norms.mean())
    
    index = faiss.IndexFlatIP(vectors.shape[1])
    index.add(vectors)
    
    faiss.write_index(index, 'faiss_index.index')

if __name__ == "__main__":
    upload()
